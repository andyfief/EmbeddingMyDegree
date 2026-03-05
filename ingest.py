#!/usr/bin/env python3
"""
ingest.py — Reads school_scan_results.json, processes every file,
and outputs chunks.jsonl ready for embedding.
No API calls — only produces chunks.jsonl and ingest_errors.log.
"""

import json
import traceback
from pathlib import Path
from collections import defaultdict

# ---------------------------------------------------------------------------
# Tunable constants
# ---------------------------------------------------------------------------

CHUNK_CHARS              = 6_000  # target characters per chunk (~1500 tokens @ 4 chars/token)
OVERLAP_CHARS            = 400   # overlap between consecutive chunks
HEADER_CHARS             = 1_000 # characters to extract for the context header
PDF_EMPTY_PAGE_THRESHOLD = 20    # min non-empty lines for a page to be valid
MIN_PDF_TOTAL_LINES      = 20    # skip whole PDF if fewer total lines
EXCEL_ROWS_PER_CHUNK     = 100   # data rows per Excel chunk

# ---------------------------------------------------------------------------
# Category mapping
# ---------------------------------------------------------------------------

CATEGORY_MAP = {
    # code / plaintext
    ".md": "code", ".mdx": "code", ".py": "code", ".cpp": "code",
    ".c": "code", ".h": "code", ".html": "code", ".js": "code",
    ".jsx": "code", ".cjs": "code", ".mjs": "code", ".css": "code",
    ".sql": "code", ".hs": "code", ".r": "code", ".bash": "code",
    ".yaml": "code", ".yml": "code", ".toml": "code", ".cfg": "code",
    ".xml": "code", ".msg": "code", ".srv": "code", ".urdf": "code",
    ".txt": "code", ".csv": "code", ".json": "code", ".ipynb": "code",
    ".asm": "code", ".sv": "code", ".sh": "code", ".tsl": "code",
    ".plain": "code", ".cmake": "code", ".in": "code",
    # documents
    ".pdf": "document", ".docx": "document",
    ".pptx": "document", ".xlsx": "document",
    # images
    ".jpg": "image", ".jpeg": "image", ".png": "image", ".webp": "image",
}

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def make_char_chunks(lines: list[str], file_path: str, ext: str) -> list[dict]:
    """Character-based chunking with context header and overlap.
    
    Produces uniform chunk sizes regardless of line length —
    critical for files where lines vary wildly (JSON, notebooks, minified code).
    Target: ~6000 chars (~1500 tokens), well under the 8192 token embedding limit.
    """
    text = "".join(lines)
    if not text.strip():
        return []

    header_text = text[:HEADER_CHARS]
    step = CHUNK_CHARS - OVERLAP_CHARS

    # Pre-compute start positions to know total_chunks upfront
    starts = list(range(0, len(text), step))
    if not starts:
        starts = [0]
    total_chunks = len(starts)

    chunks = []
    for idx, start in enumerate(starts):
        end = min(start + CHUNK_CHARS, len(text))
        body = text[start:end]
        chunk_text = f"path: {file_path}\nheader:\n{header_text}\nbody:\n{body}"
        chunks.append({
            "file_path":    file_path,
            "ext":          ext,
            "category":     CATEGORY_MAP.get(ext, "code"),
            "chunk_index":  idx,
            "total_chunks": total_chunks,
            "start_char":   start,
            "end_char":     end,
            "text":         chunk_text,
        })
    return chunks


def read_text_file(fpath: Path) -> list[str]:
    for enc in ("utf-8", "latin-1", "cp1252"):
        try:
            return fpath.read_text(encoding=enc).splitlines(keepends=True)
        except UnicodeDecodeError:
            continue
    return fpath.read_text(encoding="utf-8", errors="replace").splitlines(keepends=True)


# ---------------------------------------------------------------------------
# Per-type processors
# ---------------------------------------------------------------------------

def process_plaintext(fpath: Path, rel_path: str, ext: str) -> list[dict]:
    lines = read_text_file(fpath)
    if not lines:
        return []
    return make_char_chunks(lines, rel_path, ext)


def process_pdf(fpath: Path, rel_path: str, error_log) -> list[dict]:
    import fitz  # pymupdf

    doc = fitz.open(str(fpath))
    all_lines: list[str] = []
    chunk_page_starts: list[int] = []  # page number (1-indexed) at each line index

    skipped_pages = 0
    for page_num, page in enumerate(doc, start=1):
        text = page.get_text()
        page_lines = text.splitlines(keepends=True)
        non_empty = [l for l in page_lines if l.strip()]
        if len(non_empty) < PDF_EMPTY_PAGE_THRESHOLD:
            error_log.append(
                f"SKIP_PAGE {rel_path} page {page_num}: only {len(non_empty)} non-empty lines (image-based?)"
            )
            skipped_pages += 1
            continue
        for line in page_lines:
            all_lines.append(line)
            chunk_page_starts.append(page_num)

    doc.close()

    if len([l for l in all_lines if l.strip()]) < MIN_PDF_TOTAL_LINES:
        error_log.append(f"SKIP_FILE {rel_path}: fewer than {MIN_PDF_TOTAL_LINES} total lines")
        return []

    # Character-based chunking (consistent with all other file types)
    full_text = "".join(all_lines)
    header_text = full_text[:HEADER_CHARS]
    step = CHUNK_CHARS - OVERLAP_CHARS

    starts = list(range(0, len(full_text), step))
    if not starts:
        starts = [0]
    total_chunks = len(starts)

    chunks = []
    for idx, start in enumerate(starts):
        end = min(start + CHUNK_CHARS, len(full_text))
        body = full_text[start:end]
        chunk_text = f"path: {rel_path}\nheader:\n{header_text}\nbody:\n{body}"

        # Approximate page number: find which line index the char start falls on
        char_count = 0
        start_page = 1
        for line_idx, line in enumerate(all_lines):
            if char_count >= start:
                start_page = chunk_page_starts[line_idx] if line_idx < len(chunk_page_starts) else 1
                break
            char_count += len(line)

        chunks.append({
            "file_path":    rel_path,
            "ext":          ".pdf",
            "category":     "document",
            "chunk_index":  idx,
            "total_chunks": total_chunks,
            "start_char":   start,
            "end_char":     end,
            "start_page":   start_page,
            "text":         chunk_text,
        })
    return chunks


def process_docx(fpath: Path, rel_path: str) -> list[dict]:
    from docx import Document
    doc = Document(str(fpath))
    lines = [para.text + "\n" for para in doc.paragraphs]
    if not lines:
        return []
    return make_char_chunks(lines, rel_path, ".docx")


def process_pptx(fpath: Path, rel_path: str) -> list[dict]:
    from pptx import Presentation
    prs = Presentation(str(fpath))

    # Collect all slide text into lines, with slide markers for context
    lines = []
    for idx, slide in enumerate(prs.slides, start=1):
        lines.append(f"--- Slide {idx} ---\n")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        lines.append(t + "\n")

    # Run through make_char_chunks for consistent sizing across all file types
    return make_char_chunks(lines, rel_path, ".pptx")


def process_xlsx(fpath: Path, rel_path: str) -> list[dict]:
    import openpyxl
    wb = openpyxl.load_workbook(str(fpath), read_only=True, data_only=True)
    all_chunks = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = list(ws.iter_rows(values_only=True))
        if not rows:
            continue
        header_row = rows[0]
        header_str = ", ".join(str(c) if c is not None else "" for c in header_row)

        # Build lines: repeat the header row at the top, then all data rows
        lines = [f"Sheet: {sheet_name}\n", header_str + "\n"]
        for row in rows[1:]:
            lines.append(", ".join(str(c) if c is not None else "" for c in row) + "\n")

        # make_char_chunks handles consistent sizing
        sheet_chunks = make_char_chunks(lines, rel_path, ".xlsx")
        for chunk in sheet_chunks:
            chunk["sheet_name"] = sheet_name
        all_chunks.extend(sheet_chunks)

    wb.close()

    # Re-index chunk_index globally across all sheets
    for i, chunk in enumerate(all_chunks):
        chunk["chunk_index"]  = i
        chunk["total_chunks"] = len(all_chunks)

    return all_chunks


def process_image(fpath: Path, rel_path: str, ext: str) -> list[dict]:
    return [{
        "file_path":    rel_path,
        "ext":          ext,
        "category":     "image",
        "chunk_index":  0,
        "total_chunks": 1,
        "text":         "",
    }]


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    scan_file = Path("school_scan_results.json")
    if not scan_file.exists():
        print(f"ERROR: {scan_file} not found. Run scan_school_dir.py first.")
        return

    with open(scan_file) as f:
        scan = json.load(f)

    root = Path(scan["root"])

    # Build flat list of (abs_path, rel_path, ext)
    file_list: list[tuple[Path, str, str]] = []
    for ext, data in scan["extensions"].items():
        for example in data.get("examples", []):
            # examples are relative to School root
            pass
    # We need all files, not just examples. Reconstruct from extensions + re-walk.
    # Actually, the scan doesn't store the full list per extension beyond examples.
    # We need to re-walk using the same filters as the scan script.
    # Import scan logic or re-implement minimal version.
    # Since scan_school_dir.py is the source of truth for filters, we re-walk here
    # using the same SKIP_DIRS and SKIP_PATH_FRAGMENTS, limiting to extensions in the scan.

    ALLOWLIST = set(scan["extensions"].keys())

    SKIP_DIRS = {
        ".git", ".idea", ".vscode",
        "__pycache__", "venv", ".venv", "env", "build", "dist", "target",
        ".gradle", ".next", ".nuxt", "out", "coverage",
        ".cache", ".pytest_cache", ".mypy_cache", "eggs", ".eggs",
        "node_modules",
        ".vscode-server", ".nvm", ".dotnet", ".forever", ".local", ".ssh", ".npm",
        "sample-data",
    }
    SKIP_PATH_FRAGMENTS = {
        "site-packages", "dist-info", "scikit-learn", "pyarrow",
        "huggingface", "palm/media", "palm/tests/artifacts", "palm/lerobot",
        "cs_361/archive",
    }
    SKIP_FILENAMES = {"package-lock.json", "CMakeLists.txt"}

    import os

    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if d not in SKIP_DIRS]
        current = Path(dirpath)
        for fname in filenames:
            fpath = current / fname
            ext = fpath.suffix.lower() if fpath.suffix else ""
            if ext not in ALLOWLIST:
                continue
            if fname in SKIP_FILENAMES:
                continue
            path_str = str(fpath).replace("\\", "/").lower()
            if any(frag.lower() in path_str for frag in SKIP_PATH_FRAGMENTS):
                continue
            try:
                rel = str(fpath.relative_to(root)).replace("\\", "/")
            except ValueError:
                rel = str(fpath)
            file_list.append((fpath, rel, ext))

    total_files = len(file_list)
    print(f"Files to process: {total_files:,}")

    out_path = Path("chunks.jsonl")
    log_path = Path("ingest_errors.log")

    processed = failed = 0
    total_chunks = 0
    category_counts: dict[str, int] = defaultdict(int)
    skipped_pdf_pages = 0
    error_log_entries: list[str] = []

    with open(out_path, "w", encoding="utf-8") as out_f, \
         open(log_path, "w", encoding="utf-8") as log_f:

        for fpath, rel_path, ext in file_list:
            try:
                category = CATEGORY_MAP.get(ext, "code")
                page_errors: list[str] = []

                if category == "image":
                    chunks = process_image(fpath, rel_path, ext)

                elif ext == ".pdf":
                    chunks = process_pdf(fpath, rel_path, page_errors)
                    for e in page_errors:
                        log_f.write(e + "\n")
                        if e.startswith("SKIP_PAGE"):
                            skipped_pdf_pages += 1

                elif ext == ".docx":
                    chunks = process_docx(fpath, rel_path)

                elif ext == ".pptx":
                    chunks = process_pptx(fpath, rel_path)

                elif ext == ".xlsx":
                    chunks = process_xlsx(fpath, rel_path)

                else:
                    chunks = process_plaintext(fpath, rel_path, ext)

                for chunk in chunks:
                    out_f.write(json.dumps(chunk, ensure_ascii=False) + "\n")
                    category_counts[chunk["category"]] += 1

                total_chunks += len(chunks)
                processed += 1

            except Exception:
                failed += 1
                tb = traceback.format_exc()
                log_f.write(f"ERROR processing {rel_path}:\n{tb}\n{'-'*60}\n")

    # Summary
    print(f"\n{'='*50}")
    print(f"  Files processed : {processed:,}")
    print(f"  Files failed    : {failed:,}")
    print(f"  Total chunks    : {total_chunks:,}")
    print(f"\n  Category breakdown:")
    for cat, count in sorted(category_counts.items()):
        print(f"    {cat:<12} {count:>8,} chunks")
    print(f"\n  Skipped PDF pages (image-based): {skipped_pdf_pages:,}")
    print(f"\n  Output  -> {out_path.resolve()}")
    print(f"  Errors  -> {log_path.resolve()}")
    print(f"{'='*50}\n")


if __name__ == "__main__":
    main()